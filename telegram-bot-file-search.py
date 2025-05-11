import nest_asyncio
nest_asyncio.apply()
import os
import logging
import asyncio
import zipfile
import gc
import threading
import concurrent.futures

import docx
import openpyxl
from pptx import Presentation
import fitz  # PyMuPDF

from telegram import Update
from telegram.ext import (
    ApplicationBuilder,
    CommandHandler,
    MessageHandler,
    filters,
    ContextTypes
)

logging.basicConfig(
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
    level=logging.INFO
)
logger = logging.getLogger(__name__)

# ---------------------
# Configuration
# ---------------------
TELEGRAM_BOT_TOKEN = "YOUR_TELEGRAM_BOT_TOKEN"
SEARCH_DIRECTORY = r"C:\path\to\your\files"

AUTHORIZED_USER_IDS = {123456789, 987654321}
AUTHORIZED_USERNAMES = {"username1", "username2"}

TEXT_FILE_EXTENSIONS = {
    ".txt", ".md", ".log", ".ini", ".conf", ".cfg",
    ".py", ".js", ".ts", ".html", ".htm", ".css", ".php",
    ".c", ".cpp", ".java", ".cs", ".go", ".sh", ".rb",
    ".csv", ".json", ".xml", ".yml", ".yaml"
}
NEW_OFFICE_EXTENSIONS = {".docx", ".xlsx", ".pptx"}
OLD_OFFICE_EXTENSIONS = {".doc", ".xls", ".ppt"}
PDF_EXTENSION = {".pdf"}

MAX_MESSAGE_LENGTH = 4000

# ---------------------
# Global state
# ---------------------
files_checked = 0
files_checked_lock = threading.Lock()
# New index: file_path -> normalized content string
FILE_CONTENT_INDEX: dict[str, str] = {}

# ---------------------
# Auth
# ---------------------
def is_authorized(user: "User") -> bool:
    return (
        user.id in AUTHORIZED_USER_IDS or
        (user.username and user.username.lower() in AUTHORIZED_USERNAMES)
    )

# ---------------------
# Utils
# ---------------------
def normalize_text(text: str) -> str:
    return ' '.join(text.lower().split())

def close_all_open_files():
    gc.collect()

def is_valid_docx(file_path: str) -> bool:
    try:
        with zipfile.ZipFile(file_path, 'r') as z:
            return '[Content_Types].xml' in z.namelist()
    except Exception:
        return False

# ---------------------
# Parsers
# ---------------------
def parse_docx_file(file_path: str) -> str:
    if not is_valid_docx(file_path):
        logger.warning(f"Skipping invalid .docx: {file_path}")
        return ""
    try:
        d = docx.Document(file_path)
        return "\n".join(p.text for p in d.paragraphs)
    except Exception as e:
        logger.warning(f"docx parse error {file_path}: {e}")
        return ""

def parse_xlsx_file(file_path: str) -> str:
    try:
        wb = openpyxl.load_workbook(file_path, read_only=True)
        rows = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    rows.append(" ".join(cells))
        return "\n".join(rows)
    except Exception as e:
        logger.warning(f"xlsx parse error {file_path}: {e}")
        return ""

def parse_pptx_file(file_path: str) -> str:
    try:
        pres = Presentation(file_path)
        texts = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception as e:
        logger.warning(f"pptx parse error {file_path}: {e}")
        return ""

def parse_text_file(file_path: str) -> str:
    try:
        with open(file_path, "r", errors="ignore") as f:
            return f.read()
    except Exception as e:
        logger.warning(f"text file error {file_path}: {e}")
        return ""

def parse_pdf_file(file_path: str) -> str:
    try:
        doc = fitz.open(file_path)
        return "".join(page.get_text() for page in doc)
    except Exception as e:
        logger.warning(f"PDF parse error {file_path}: {e}")
        return ""

def extract_file_content(file_path: str) -> str:
    _, ext = os.path.splitext(file_path.lower())
    if os.path.basename(file_path).startswith("~$"):
        return ""
    if ext in OLD_OFFICE_EXTENSIONS:
        return ""
    if ext in TEXT_FILE_EXTENSIONS:
        return parse_text_file(file_path)
    if ext == ".docx":
        return parse_docx_file(file_path)
    if ext == ".xlsx":
        return parse_xlsx_file(file_path)
    if ext == ".pptx":
        return parse_pptx_file(file_path)
    if ext in PDF_EXTENSION:
        return parse_pdf_file(file_path)
    return ""

# ---------------------
# Build exact-string index
# ---------------------
def index_file_content(path: str) -> tuple[str,str]:
    text = extract_file_content(path)
    norm = normalize_text(text) if text else ""
    return path, norm

def build_content_index(search_dir: str):
    global FILE_CONTENT_INDEX, files_checked
    files = []
    for root, _, fnames in os.walk(search_dir):
        for f in fnames:
            if not f.startswith("~$"):
                files.append(os.path.join(root, f))
    local_idx: dict[str,str] = {}
    with concurrent.futures.ThreadPoolExecutor(max_workers=32) as pool:
        futures = {pool.submit(index_file_content, p): p for p in files}
        for fut in concurrent.futures.as_completed(futures):
            path, content = fut.result()
            with files_checked_lock:
                files_checked += 1
            if content:
                local_idx[path] = content
    FILE_CONTENT_INDEX = local_idx

def search_exact_index(query: str) -> list[str]:
    nq = normalize_text(query)
    return [p for p, content in FILE_CONTENT_INDEX.items() if nq in content]

# ---------------------
# Fallback live search
# ---------------------
def process_file_for_search(path: str, nq: str):
    global files_checked
    with files_checked_lock:
        files_checked += 1
    text = extract_file_content(path)
    return path if text and nq in normalize_text(text) else None

def find_string_in_files_parallel(search_dir: str, query: str):
    all_paths = [
        os.path.join(root, f)
        for root, _, fnames in os.walk(search_dir)
        for f in fnames
        if not f.startswith("~$")
    ]
    nq = normalize_text(query)
    results = []
    with concurrent.futures.ThreadPoolExecutor(max_workers=32) as pool:
        futs = {pool.submit(process_file_for_search, p, nq): p for p in all_paths}
        for fut in concurrent.futures.as_completed(futs):
            r = fut.result()
            if r:
                results.append(r)
    return results

def chunk_text(t: str, max_size=MAX_MESSAGE_LENGTH):
    return [t[i:i+max_size] for i in range(0, len(t), max_size)]

# ---------------------
# Bot handlers
# ---------------------
async def start_command(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    if not is_authorized(update.effective_user):
        return await update.message.reply_text("You are not authorized.")
    await update.message.reply_text(
        "Use /index to build the exact-string index.\n"
        "Then send any text and I'll return files containing that exact string."
    )

async def index_command(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    if not is_authorized(update.effective_user):
        return await update.message.reply_text("You are not authorized.")
    global files_checked
    with files_checked_lock:
        files_checked = 0

    msg = await update.message.reply_text("Indexing files: 0")
    done = False

    async def progress_updater():
        while not done:
            await asyncio.sleep(1)
            with files_checked_lock:
                n = files_checked
            try:
                await ctx.bot.edit_message_text(
                    chat_id=update.effective_chat.id,
                    message_id=msg.message_id,
                    text=f"Indexing files: {n}"
                )
            except Exception as e:
                if "Message is not modified" not in str(e):
                    logger.warning(f"Progress update failed: {e}")

    task = asyncio.create_task(progress_updater())
    await asyncio.to_thread(build_content_index, SEARCH_DIRECTORY)
    done = True
    await task

    await update.message.reply_text(
        f"Index built: {len(FILE_CONTENT_INDEX)} files indexed."
    )
    close_all_open_files()

async def handle_text(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    if not is_authorized(update.effective_user):
        return await update.message.reply_text("You are not authorized.")

    query = update.message.text.strip()
    if not query:
        return await update.message.reply_text("Please send a non-empty query.")

    # use index if available
    if FILE_CONTENT_INDEX:
        results = search_exact_index(query)
    else:
        global files_checked
        with files_checked_lock:
            files_checked = 0
        msg = await update.message.reply_text("Files checked: 0")
        done = False

        async def progress_updater():
            while not done:
                await asyncio.sleep(1)
                with files_checked_lock:
                    n = files_checked
                try:
                    await ctx.bot.edit_message_text(
                        chat_id=update.effective_chat.id,
                        message_id=msg.message_id,
                        text=f"Files checked: {n}"
                    )
                except Exception as e:
                    if "Message is not modified" not in str(e):
                        logger.warning(f"Progress update failed: {e}")

        task = asyncio.create_task(progress_updater())
        results = await asyncio.to_thread(find_string_in_files_parallel, SEARCH_DIRECTORY, query)
        done = True
        await task

    total = len(results)
    if total:
        out = "Found in:\n" + "\n".join(results)
    else:
        out = f"No files contain exactly '{query}'."

    for chunk in chunk_text(out):
        await update.message.reply_text(chunk)
    await update.message.reply_text(f"Search done. {total} result(s).")
    close_all_open_files()

async def main():
    app = ApplicationBuilder().token(TELEGRAM_BOT_TOKEN).build()
    app.add_handler(CommandHandler("start", start_command))
    app.add_handler(CommandHandler("index", index_command))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_text))
    await app.run_polling()

if __name__ == "__main__":
    asyncio.run(main())
